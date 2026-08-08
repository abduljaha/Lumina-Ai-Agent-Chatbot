"""Document processing - loading, chunking, and metadata extraction.

Every extractor returns `(text, metadata)` - `metadata` carries document-level
facts (extraction method, page/slide counts, OCR usage) that get merged into
every chunk's own metadata in `process()`, so a caller can tell *how* a
document's content was obtained (native text vs. OCR vs. a vision-LLM
description) without re-deriving it, and a genuine failure (password-
protected PDF, corrupt file) carries a specific reason rather than a bare
empty string indistinguishable from "this file has no content".
"""
from __future__ import annotations

import asyncio
import logging
from pathlib import Path
from typing import Any

logger = logging.getLogger("app")

# File extensions this processor can turn into indexable text. Shared with
# the upload endpoint's allow-list (see files.py) so the two can't drift
# apart - an extension accepted there but not handled here would silently
# index zero content, and one handled here but rejected there would just be
# dead code.
CODE_EXTENSIONS = {
    "py", "js", "jsx", "ts", "tsx", "java", "c", "cpp", "h", "hpp", "cs",
    "go", "rs", "rb", "php", "swift", "kt", "sql", "sh", "yaml", "yml",
    "toml", "ini", "css", "scss",
}
IMAGE_EXTENSIONS = {"jpg", "jpeg", "png", "webp", "bmp", "tiff", "tif"}
SUPPORTED_EXTENSIONS = (
    {"pdf", "docx", "pptx", "txt", "md", "csv", "xlsx", "xls", "json", "xml", "html", "htm"}
    | CODE_EXTENSIONS
    | IMAGE_EXTENSIONS
)

# Bounds how many embedded images per document get OCR'd/described - a
# pathological file with hundreds of images would otherwise turn one upload
# into hundreds of OCR calls (and, for the vision-description fallback,
# hundreds of LLM calls) before it ever finishes processing.
_MAX_EMBEDDED_IMAGES = 20
# Below this many characters of *native* text, a PDF page is treated as
# scanned/image-based rather than a real (if short) text layer - real pages
# almost never extract to just a handful of stray characters, but a
# scanned page's native layer is usually completely empty or near-empty
# OCR artifacts left by the scanner software.
_NATIVE_TEXT_MIN_CHARS = 20

# ---------------------------------------------------------------------------
# Shared OCR backend - module-level so the (expensive to construct) engine
# is loaded once per process and reused across every DocumentProcessor
# instance and every upload, not reloaded on every single call.
# ---------------------------------------------------------------------------
_rapidocr_engine: Any = None
_rapidocr_load_attempted = False


async def _get_rapidocr_engine() -> Any:
    """Lazily construct and cache the RapidOCR engine (pure Python, no system binary)."""
    global _rapidocr_engine, _rapidocr_load_attempted
    if _rapidocr_load_attempted:
        return _rapidocr_engine
    _rapidocr_load_attempted = True
    try:
        from rapidocr_onnxruntime import RapidOCR

        _rapidocr_engine = await asyncio.to_thread(RapidOCR)
        logger.info("RapidOCR engine loaded")
    except Exception as exc:  # noqa: BLE001
        logger.warning("RapidOCR unavailable (%s) - will try pytesseract if installed", exc)
        _rapidocr_engine = None
    return _rapidocr_engine


async def ocr_image(image: Any) -> str:
    """Run OCR on a PIL Image, preferring RapidOCR, falling back to pytesseract.

    Two independent OCR backends rather than one: RapidOCR needs no system
    binary (reliable in any environment this app runs in, including ones
    without admin rights to install Tesseract), but pytesseract is tried
    second in case Tesseract *is* present and RapidOCR's models fail to
    load for some reason - either succeeding is enough.
    """
    engine = await _get_rapidocr_engine()
    if engine is not None:
        try:
            import numpy as np

            def _run() -> str:
                result, _ = engine(np.array(image))
                if not result:
                    return ""
                return "\n".join(line[1] for line in result)

            text = await asyncio.to_thread(_run)
            if text.strip():
                return text
        except Exception as exc:  # noqa: BLE001
            logger.warning("RapidOCR failed on an image, trying pytesseract: %s", exc)

    try:
        import pytesseract

        text = await asyncio.to_thread(pytesseract.image_to_string, image)
        return text
    except Exception as exc:  # noqa: BLE001
        logger.warning("pytesseract also unavailable/failed: %s", exc)
        return ""


class DocumentProcessor:
    """Loads, chunks, and extracts metadata from uploaded documents.

    Supported types: PDF (native text, with automatic OCR fallback per page
    for scanned/image-based PDFs), DOCX/PPTX (paragraph/slide text plus OCR
    of embedded images), TXT/Markdown, CSV, Excel (XLSX/XLS), JSON, XML,
    HTML, source code files, and standalone images (OCR, with a vision-LLM
    description as a further fallback for images with little/no text - a
    photo or diagram, not a scanned document). Anything not explicitly
    handled below still falls through to a plain-text read, so it degrades
    rather than hard-fails.
    """

    CHUNK_SIZE = 800
    CHUNK_OVERLAP = 100

    def __init__(self, chunk_size: int = CHUNK_SIZE, chunk_overlap: int = CHUNK_OVERLAP) -> None:
        self._splitter = None
        try:
            from langchain_text_splitters import RecursiveCharacterTextSplitter

            self._splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                separators=["\n\n", "\n", ". ", " ", ""],
            )
        except Exception as exc:  # noqa: BLE001
            logger.warning(
                "langchain_text_splitters not available; using fallback splitter: %s",
                exc,
            )

    async def process(self, file_path: str, filename: str) -> tuple[list[dict[str, Any]], dict[str, Any]]:
        """Process a file, returning (chunks, document-level metadata).

        `metadata` always carries `extraction_method` on success, or a
        specific, human-readable `error` when `chunks` comes back empty -
        callers (see files.py's upload endpoint) surface that reason
        directly instead of a single generic "no text found" message that
        can't distinguish a password-protected PDF from a genuinely blank
        file from a missing OCR engine.
        """
        file_ext = Path(filename).suffix.lower().lstrip(".")
        text, doc_metadata = await self._extract_text(file_path, file_ext)

        if not text.strip():
            reason = doc_metadata.get("error") or (
                "No text could be extracted from this file. It may be empty, corrupted, "
                "or in an unexpected format."
            )
            logger.warning("No text extracted from %s (%s)", filename, reason)
            return [], {**doc_metadata, "error": reason}

        if self._splitter is not None:
            chunks = self._splitter.split_text(text)
        else:
            chunks = self._fallback_split_text(text)
        results = []
        for i, chunk in enumerate(chunks):
            metadata = self._extract_metadata(chunk, filename, i, file_ext)
            metadata.update(doc_metadata)
            results.append(
                {
                    "content": chunk,
                    "chunk_index": i,
                    "source": filename,
                    "metadata": metadata,
                }
            )
        logger.info(
            "Processed %s into %d chunk(s) via %s",
            filename, len(chunks), doc_metadata.get("extraction_method", "native"),
        )
        return results, doc_metadata

    async def _extract_text(self, file_path: str, file_ext: str) -> tuple[str, dict[str, Any]]:
        """Extract text from a file based on its extension."""
        if file_ext == "pdf":
            return await self._extract_pdf(file_path)
        if file_ext == "docx":
            return await self._extract_docx(file_path)
        if file_ext == "pptx":
            return await self._extract_pptx(file_path)
        if file_ext in ("txt", "md", "csv", "json") or file_ext in CODE_EXTENSIONS:
            # JSON and source files are already human-readable plain text -
            # no structural parsing needed to make them embeddable/searchable.
            return await self._extract_text_simple(file_path)
        if file_ext in ("xlsx", "xls"):
            return await self._extract_excel(file_path)
        if file_ext == "xml":
            return await self._extract_xml(file_path)
        if file_ext in ("html", "htm"):
            return await self._extract_html(file_path)
        if file_ext in IMAGE_EXTENSIONS:
            return await self._extract_image(file_path)
        if file_ext in ("zip", "mp3", "wav", "mp4"):
            return await self._extract_unsupported(file_ext)
        return await self._extract_text_simple(file_path)

    async def _extract_pdf(self, file_path: str) -> tuple[str, dict[str, Any]]:
        """Extract text from a PDF, OCR-ing any page whose native text layer is empty/near-empty."""
        try:
            from pypdf import PdfReader

            try:
                reader = PdfReader(file_path)
            except Exception as exc:  # noqa: BLE001
                # pypdf raises different exception types across versions for
                # a malformed vs. an encrypted PDF - the message text is the
                # one reliable signal for telling a user "this needs a
                # password" apart from "this file is corrupt".
                if "password" in str(exc).lower() or "encrypt" in str(exc).lower():
                    return "", {"error": "This PDF is password-protected and can't be read."}
                return "", {"error": f"Could not open PDF: {exc}"}

            if reader.is_encrypted:
                try:
                    reader.decrypt("")  # some "restricted" (not password-protected) PDFs open with an empty password
                except Exception:  # noqa: BLE001
                    pass
                if reader.is_encrypted:
                    return "", {"error": "This PDF is password-protected and can't be read."}

            page_texts: list[str] = []
            ocr_page_count = 0
            for i, page in enumerate(reader.pages):
                native = (page.extract_text() or "").strip()
                if len(native) >= _NATIVE_TEXT_MIN_CHARS:
                    page_texts.append(native)
                    continue
                # No real text layer on this page - likely a scanned image,
                # so render it and OCR it instead of just skipping it.
                ocr_text = (await self._ocr_pdf_page(file_path, i)).strip()
                if ocr_text:
                    ocr_page_count += 1
                    page_texts.append(ocr_text)
                elif native:
                    page_texts.append(native)

            page_count = len(reader.pages)
            if ocr_page_count == 0:
                method = "native"
            elif ocr_page_count == page_count:
                method = "ocr"
            else:
                method = "hybrid"

            return "\n\n".join(page_texts), {
                "extraction_method": method,
                "page_count": page_count,
                "ocr_page_count": ocr_page_count,
            }
        except Exception as exc:  # noqa: BLE001
            logger.error("PDF extraction failed: %s", exc)
            return "", {"error": str(exc)}

    async def _ocr_pdf_page(self, file_path: str, page_index: int) -> str:
        """Render one PDF page to an image (via PyMuPDF) and OCR it."""
        try:
            import fitz  # PyMuPDF - pure wheel, no system binary needed

            def _render() -> Any:
                from PIL import Image

                doc = fitz.open(file_path)
                try:
                    page = doc[page_index]
                    # 2x zoom: typical scan resolutions OCR noticeably more
                    # accurately above the PDF's default ~72 DPI render.
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    return Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                finally:
                    doc.close()

            image = await asyncio.to_thread(_render)
            return await ocr_image(image)
        except ImportError:
            logger.warning("PyMuPDF not installed - cannot OCR scanned PDF pages")
            return ""
        except Exception as exc:  # noqa: BLE001
            logger.warning("OCR of PDF page %d failed: %s", page_index, exc)
            return ""

    async def _extract_docx(self, file_path: str) -> tuple[str, dict[str, Any]]:
        """Extract paragraph text from a DOCX file, plus OCR text from any embedded images."""
        try:
            from docx import Document

            doc = Document(file_path)
            paragraph_text = "\n".join(p.text for p in doc.paragraphs if p.text)

            image_parts = [rel.target_part for rel in doc.part.rels.values() if "image" in rel.reltype]
            image_texts = await self._ocr_embedded_images(
                (part.blob for part in image_parts[:_MAX_EMBEDDED_IMAGES]), file_path
            )

            combined = "\n\n".join(p for p in [paragraph_text, "\n\n".join(image_texts)] if p)
            method = "native"
            if image_texts:
                method = "hybrid" if paragraph_text.strip() else "ocr"
            return combined, {
                "extraction_method": method,
                "embedded_image_count": len(image_parts),
                "embedded_images_with_text": len(image_texts),
            }
        except Exception as exc:  # noqa: BLE001
            logger.error("DOCX extraction failed: %s", exc)
            return "", {"error": str(exc)}

    async def _extract_pptx(self, file_path: str) -> tuple[str, dict[str, Any]]:
        """Extract text from a PowerPoint file, slide by slide, plus OCR of any picture shapes."""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(file_path)
            slides: list[str] = []
            image_blobs: list[bytes] = []
            slide_pictures: list[tuple[int, list[Any]]] = []

            for i, slide in enumerate(prs.slides):
                texts = [
                    shape.text_frame.text
                    for shape in slide.shapes
                    if shape.has_text_frame and shape.text_frame.text.strip()
                ]
                pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
                slide_pictures.append((i, texts))
                for shape in pictures[: max(0, _MAX_EMBEDDED_IMAGES - len(image_blobs))]:
                    image_blobs.append(shape.image.blob)

            image_texts = await self._ocr_embedded_images(iter(image_blobs), file_path)
            # Embedded-image text isn't tied back to a specific slide (OCR
            # runs on the flattened image list above) - appended as its own
            # section rather than guessed into a particular slide's text.
            for i, texts in slide_pictures:
                if texts:
                    slides.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
            if image_texts:
                slides.append("[Embedded images]\n" + "\n\n".join(image_texts))

            method = "native"
            if image_texts:
                method = "hybrid" if any(t for _, t in slide_pictures) else "ocr"
            return "\n\n".join(slides), {
                "extraction_method": method,
                "slide_count": len(prs.slides),
                "embedded_image_count": len(image_blobs),
                "embedded_images_with_text": len(image_texts),
            }
        except Exception as exc:  # noqa: BLE001
            logger.error("PPTX extraction failed: %s", exc)
            return "", {"error": str(exc)}

    async def _ocr_embedded_images(self, blobs: Any, source_file: str) -> list[str]:
        """OCR a sequence of raw image blobs (from DOCX/PPTX), returning non-empty results."""
        from PIL import Image
        import io

        texts: list[str] = []
        for i, blob in enumerate(blobs):
            try:
                image = await asyncio.to_thread(Image.open, io.BytesIO(blob))
                image.load()
                text = (await ocr_image(image)).strip()
                if text:
                    texts.append(f"[Image {i + 1}]: {text}")
            except Exception as exc:  # noqa: BLE001
                logger.warning("Could not OCR embedded image %d in %s: %s", i + 1, source_file, exc)
        return texts

    async def _extract_xml(self, file_path: str) -> tuple[str, dict[str, Any]]:
        """Extract visible text content from an XML file, tags stripped."""
        try:
            from xml.etree import ElementTree

            tree = ElementTree.parse(file_path)
            parts = [node.strip() for node in tree.getroot().itertext() if node and node.strip()]
            return "\n".join(parts), {"extraction_method": "native"}
        except Exception as exc:  # noqa: BLE001
            logger.error("XML extraction failed: %s", exc)
            return "", {"error": str(exc)}

    async def _extract_html(self, file_path: str) -> tuple[str, dict[str, Any]]:
        """Extract visible text content from an HTML file, tags/scripts/styles stripped."""
        try:
            from bs4 import BeautifulSoup

            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            return soup.get_text(separator="\n", strip=True), {"extraction_method": "native"}
        except Exception as exc:  # noqa: BLE001
            logger.error("HTML extraction failed: %s", exc)
            return "", {"error": str(exc)}

    async def _extract_text_simple(self, file_path: str) -> tuple[str, dict[str, Any]]:
        """Extract text from a simple text file."""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(), {"extraction_method": "native"}
        except Exception as exc:  # noqa: BLE001
            logger.error("Text extraction failed: %s", exc)
            return "", {"error": str(exc)}

    async def _extract_excel(self, file_path: str) -> tuple[str, dict[str, Any]]:
        """Extract text from an Excel file."""
        try:
            import pandas as pd

            dfs = pd.read_excel(file_path, sheet_name=None)
            parts = []
            for sheet_name, df in dfs.items():
                parts.append(f"Sheet: {sheet_name}\n{df.to_string()}")
            return "\n\n".join(parts), {"extraction_method": "native", "sheet_count": len(dfs)}
        except Exception as exc:  # noqa: BLE001
            logger.error("Excel extraction failed: %s", exc)
            return "", {"error": str(exc)}

    async def _extract_image(self, file_path: str) -> tuple[str, dict[str, Any]]:
        """OCR a standalone image, falling back to a vision-LLM description if OCR finds little/no text.

        OCR alone only helps for images that *contain* text (a screenshot, a
        scanned page saved as a photo). A genuine photo or diagram has
        nothing for OCR to find, which used to mean it was uploaded
        successfully but indexed with zero content - permanently
        unfindable and unanswerable. The vision-description fallback covers
        that case with an actual description of what's in the image.
        """
        try:
            from PIL import Image

            image = await asyncio.to_thread(Image.open, file_path)
            image.load()
            ocr_text = (await ocr_image(image)).strip()
            if len(ocr_text) >= _NATIVE_TEXT_MIN_CHARS:
                return ocr_text, {"extraction_method": "ocr"}

            description = await self._describe_image_with_vision(file_path)
            if description:
                combined = "\n\n".join(p for p in [ocr_text, description] if p)
                return combined, {"extraction_method": "vision_description"}

            if ocr_text:
                return ocr_text, {"extraction_method": "ocr"}
            return "", {
                "error": "No text found in this image, and no vision-capable AI provider "
                "is currently available to describe it."
            }
        except Exception as exc:  # noqa: BLE001
            logger.error("Image extraction failed: %s", exc)
            return "", {"error": str(exc)}

    async def _describe_image_with_vision(self, file_path: str) -> str:
        """Ask a vision-capable LLM to describe an image, for search indexing."""
        try:
            import base64

            from app.core.container import app_container
            from app.llm.base import LLMRequest

            with open(file_path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode()
            ext = Path(file_path).suffix.lstrip(".").lower()
            mime = {"jpg": "jpeg"}.get(ext, ext) or "png"
            data_url = f"data:image/{mime};base64,{b64}"

            request = LLMRequest(
                messages=[
                    {
                        "role": "system",
                        "content": "Describe this image in detail for search indexing: objects, "
                        "any visible text, colors, layout, and overall context. Be specific and factual.",
                    },
                    {"role": "user", "content": "Describe this image."},
                ],
                images=[data_url],
                max_tokens=400,
            )
            response = await app_container.model_router.generate_vision(request)
            return (response.content or "").strip()
        except Exception as exc:  # noqa: BLE001
            logger.warning("Vision description fallback unavailable for %s: %s", file_path, exc)
            return ""

    async def _extract_unsupported(self, file_ext: str) -> tuple[str, dict[str, Any]]:
        """Handle unsupported file types gracefully."""
        return "", {"error": f"'.{file_ext}' files are not supported for indexing."}

    def _fallback_split_text(self, text: str) -> list[str]:
        """Fallback text splitter when langchain_text_splitters is unavailable."""
        if not text:
            return []
        sentences = text.split("\n\n")
        chunks: list[str] = []
        current = []
        current_len = 0
        for paragraph in sentences:
            paragraph = paragraph.strip()
            if not paragraph:
                continue
            paragraph_len = len(paragraph)
            if current_len + paragraph_len + 1 > self.CHUNK_SIZE and current:
                chunks.append("\n\n".join(current))
                current = [paragraph]
                current_len = paragraph_len
            else:
                current.append(paragraph)
                current_len += paragraph_len + 2
        if current:
            chunks.append("\n\n".join(current))
        return chunks

    def _extract_metadata(self, chunk: str, filename: str, index: int, file_type: str = "") -> dict[str, Any]:
        """Extract per-chunk metadata.

        `file_type` is stored so retrieval can filter by it (e.g. "only
        search my spreadsheets") - see Retriever.retrieve's `filters` param.
        """
        words = len(chunk.split())
        return {
            "chunk_index": index,
            "filename": filename,
            "file_type": file_type,
            "word_count": words,
            "char_count": len(chunk),
        }
